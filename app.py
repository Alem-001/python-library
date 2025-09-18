from fastapi import FastAPI, UploadFile, Form
from fastapi.responses import StreamingResponse, JSONResponse
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
import io, json

app = FastAPI()

# -------------------------------
# Helpers
# -------------------------------

def _replace_in_textframe(tf, mapping: dict):
    """Replace {{tokens}} in a full text frame (handles split runs)."""
    text = tf.text or ""
    for k, v in mapping.items():
        if isinstance(v, (dict, list)):
            continue
        text = text.replace(f"{{{{{k}}}}}", str(v or ""))
    tf.text = text

def _replace_in_table_cells(table, mapping: dict):
    """Replace {{tokens}} inside an existing table."""
    for row in table.rows:
        for cell in row.cells:
            if cell.text and "{{" in cell.text:
                new_text = cell.text
                for k, v in mapping.items():
                    if isinstance(v, (dict, list)):
                        continue
                    new_text = new_text.replace(f"{{{{{k}}}}}", str(v or ""))
                if new_text != cell.text:
                    cell.text = new_text

def _remove_shape(slide, shape):
    sp = shape._element
    sp.getparent().remove(sp)

def _normalize_table_font(table, size_pt=12):
    for row in table.rows:
        for cell in row.cells:
            if cell.text_frame:
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(size_pt)

# ---------- NEW: generic quality table (any headers, all stringy values)
def _rebuild_simple_table(slide, placeholder_shape, table_spec: dict):
    """
    Build a generic table for quality data.
    Expected JSON:
    {
      "headers": ["Label", "Value"],      # optional; default = ["Label", "Value"]
      "rows": [
        ["Business at a glance", "Flat-roof specialist"],
        ["HQ", "Bad Kreuznach, DE"],
        ["Employees", "19 + 2 MDs"],
        ["Reason to sell", "Succession (age)"]
      ]
    }
    Notes:
    - All values are treated as strings (even if they look like numbers).
    - Number of columns = len(headers).
    - Each row must have same length as headers.
    """
    headers = table_spec.get("headers") or ["Label", "Value"]
    rows = table_spec.get("rows") or []

    n_rows = 1 + len(rows)
    n_cols = len(headers)

    left, top, width, height = placeholder_shape.left, placeholder_shape.top, placeholder_shape.width, placeholder_shape.height
    _remove_shape(slide, placeholder_shape)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # Header
    for c, h in enumerate(headers):
        table.cell(0, c).text = str(h)

    # Data
    for r, rowvals in enumerate(rows, start=1):
        for c, val in enumerate(rowvals[:n_cols]):
            table.cell(r, c).text = str(val)

    _normalize_table_font(table, 12)

# ---------- Financial table (kept here for later use if needed)
def _rebuild_table_from_json(slide, placeholder_shape, financials: dict):
    years = financials.get("years", [])
    rows = financials.get("rows", [])

    n_rows = 1 + len(rows)
    n_cols = 1 + len(years)

    left, top, width, height = placeholder_shape.left, placeholder_shape.top, placeholder_shape.width, placeholder_shape.height
    _remove_shape(slide, placeholder_shape)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    table.cell(0, 0).text = "Item"
    for c, y in enumerate(years, start=1):
        if c < len(table.columns):
            table.cell(0, c).text = str(y)

    for r, rowdata in enumerate(rows, start=1):
        table.cell(r, 0).text = str(rowdata.get("item", ""))
        vals = rowdata.get("values", [])
        for c, val in enumerate(vals, start=1):
            if c < len(table.columns):
                table.cell(r, c).text = str(val)

    _normalize_table_font(table, 12)

def _walk_shapes(container, mapping: dict):
    # container can be slide or a group shape
    for shp in container.shapes:
        # Build table from marker like {{TABLE:quality}} or {{TABLE:financials}}
        if getattr(shp, "has_text_frame", False) and shp.text and "{{TABLE:" in shp.text:
            key = shp.text.strip().replace("{{TABLE:", "").replace("}}", "").strip()
            if key in mapping and isinstance(mapping[key], dict):
                spec = mapping[key]
                # If caller passed "headers"/"rows" => simple quality table
                if "rows" in spec and isinstance(spec["rows"], list):
                    _rebuild_simple_table(container, shp, spec)
                # If caller passed "years"/"rows" => financial table
                elif "years" in spec and "rows" in spec:
                    _rebuild_table_from_json(container, shp, spec)
            continue

        # Normal text boxes
        if getattr(shp, "has_text_frame", False):
            _replace_in_textframe(shp.text_frame, mapping)

        # Tables with token placeholders
        if getattr(shp, "has_table", False):
            _replace_in_table_cells(shp.table, mapping)

        # Groups (recurse)
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            _walk_shapes(shp, mapping)

# -------------------------------
# Routes
# -------------------------------

@app.get("/health")
def health():
    return {"ok": True}

@app.post("/fill")
async def fill(
    template: UploadFile,
    data: UploadFile | None = None,
    json_text: str = Form(None),
):
    try:
        if data:
            mapping = json.loads((await data.read()).decode("utf-8"))
        elif json_text:
            mapping = json.loads(json_text)
        else:
            return JSONResponse({"error": "No JSON provided"}, status_code=400)

        prs = Presentation(io.BytesIO(await template.read()))

        for slide in prs.slides:
            _walk_shapes(slide, mapping)

        out = io.BytesIO()
        prs.save(out)
        out.seek(0)
        return StreamingResponse(
            out,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": 'attachment; filename="filled.pptx"'},
        )
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)
